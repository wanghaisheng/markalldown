from flask import Flask, request, jsonify, render_template
import os
import urllib.request
from io import BytesIO
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import zipfile
import io

app = Flask(__name__)

# Function to extract text from various file formats
def extract_text_from_file(file_path, file_extension):
    text = ""
    if file_extension == '.pdf':
        text = extract_text(file_path)  # PDF file
    elif file_extension == '.docx':
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])  # DOCX file
    elif file_extension == '.pptx':
        presentation = Presentation(file_path)
        text = "\n".join([slide.shapes[0].text for slide in presentation.slides if slide.shapes])
    elif file_extension == '.txt':
        with open(file_path, 'r') as file:
            text = file.read()
    elif file_extension in ['.jpg', '.jpeg', '.png']:
        from PIL import Image
        import pytesseract
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
    return text

# Route to serve the index.html
@app.route('/')
def home():
    return render_template('index.html')

# Route to handle file upload for conversion
@app.route('/convert', methods=['POST'])
def convert():
    data = request.get_json()

    if 'url' in data:
        url = data['url']
        try:
            # Handle URL download
            response = urllib.request.urlopen(url)
            file_data = response.read()
            file_name = url.split("/")[-1]
            file_extension = os.path.splitext(file_name)[1].lower()

            # Create a temporary file from URL data
            temp_file = BytesIO(file_data)
            text_content = extract_text_from_file(temp_file, file_extension)
            return jsonify({'text_content': text_content})
        except Exception as e:
            return jsonify({'error': str(e)}), 400
    else:
        # Handle file upload
        file = request.files['file']
        if file:
            file_name = file.filename
            file_extension = os.path.splitext(file_name)[1].lower()

            # Save file temporarily and process it
            temp_file = BytesIO(file.read())
            text_content = extract_text_from_file(temp_file, file_extension)

            return jsonify({'text_content': text_content})
        else:
            return jsonify({'error': 'No file uploaded'}), 400

if __name__ == '__main__':
    app.run(debug=True)
